import os
from typing import List, Optional
from lxml import etree
from services.presentation.services.html_to_text_runs_service import (
    parse_html_text_to_text_runs as parse_inline_html_to_runs,
)

from pptx import Presentation
from pptx.shapes.autoshape import Shape
from pptx.slide import Slide
from pptx.text.text import _Paragraph, TextFrame, Font, _Run
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from lxml.etree import fromstring, tostring
from PIL import Image
from pptx.oxml.xmlchemy import OxmlElement

from pptx.util import Pt
from pptx.dml.color import RGBColor

from services.presentation.models.pptx_models import (
    PptxAutoShapeBoxModel,
    PptxBoxShapeEnum,
    PptxConnectorModel,
    PptxFillModel,
    PptxFontModel,
    PptxParagraphModel,
    PptxPictureBoxModel,
    PptxPositionModel,
    PptxPresentationModel,
    PptxShadowModel,
    PptxSlideModel,
    PptxSpacingModel,
    PptxStrokeModel,
    PptxTextBoxModel,
    PptxTextRunModel,
)
from services.presentation.utils.download_helpers import download_files
from services.presentation.utils.image_utils import (
    clip_image,
    create_circle_image,
    fit_image,
    invert_image,
    round_image_corners,
    set_image_opacity,
)
import uuid

BLANK_SLIDE_LAYOUT = 6


class PptxPresentationCreator:

    def __init__(self, ppt_model: PptxPresentationModel, temp_dir: str):
        self._temp_dir = temp_dir

        self._ppt_model = ppt_model
        self._slide_models = ppt_model.slides

        self._ppt = Presentation()
        self._ppt.slide_width = Pt(1280)
        self._ppt.slide_height = Pt(720)

    def get_sub_element(self, parent, tagname, **kwargs):
        """Helper method to create XML elements"""
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

    async def fetch_network_assets(self):
        image_urls = []
        models_with_network_asset: List[PptxPictureBoxModel] = []

        if self._ppt_model.shapes:
            for each_shape in self._ppt_model.shapes:
                if isinstance(each_shape, PptxPictureBoxModel):
                    image_path = each_shape.picture.path
                    if image_path.startswith("http"):
                        if "app_data/" in image_path:
                            relative_path = image_path.split("app_data/")[1]
                            each_shape.picture.path = os.path.join(
                                "/app_data", relative_path
                            )
                            each_shape.picture.is_network = False
                            continue
                        image_urls.append(image_path)
                        models_with_network_asset.append(each_shape)

        for each_slide in self._slide_models:
            for each_shape in each_slide.shapes:
                if isinstance(each_shape, PptxPictureBoxModel):
                    image_path = each_shape.picture.path
                    if image_path.startswith("http"):
                        if "app_data" in image_path:
                            relative_path = image_path.split("app_data/")[1]
                            each_shape.picture.path = os.path.join(
                                "/app_data", relative_path
                            )
                            each_shape.picture.is_network = False
                            continue
                        image_urls.append(image_path)
                        models_with_network_asset.append(each_shape)

        if image_urls:
            image_paths = await download_files(image_urls, self._temp_dir)

            for each_shape, each_image_path in zip(
                models_with_network_asset, image_paths
            ):
                if each_image_path:
                    each_shape.picture.path = each_image_path
                    each_shape.picture.is_network = False

    async def create_ppt(self):
        await self.fetch_network_assets()

        for slide_model in self._slide_models:
            # Adding global shapes to slide
            if self._ppt_model.shapes:
                slide_model.shapes.append(self._ppt_model.shapes)

            self.add_and_populate_slide(slide_model)

    def set_presentation_theme(self):
        slide_master = self._ppt.slide_master
        slide_master_part = slide_master.part

        theme_part = slide_master_part.part_related_by(RT.THEME)
        theme = fromstring(theme_part.blob)

        theme_colors = self._theme.colors.theme_color_mapping
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

        for color_name, hex_value in theme_colors.items():
            if color_name:
                color_element = theme.xpath(
                    f"a:themeElements/a:clrScheme/a:{color_name}/a:srgbClr",
                    namespaces=nsmap,
                )[0]
                color_element.set("val", hex_value.encode("utf-8"))

        theme_part._blob = tostring(theme)

    def add_and_populate_slide(self, slide_model: PptxSlideModel):
        slide = self._ppt.slides.add_slide(self._ppt.slide_layouts[BLANK_SLIDE_LAYOUT])

        if slide_model.background:
            self.apply_fill_to_shape(slide.background, slide_model.background)

        if slide_model.note:
            slide.notes_slide.notes_text_frame.text = slide_model.note

        for shape_model in slide_model.shapes:
            model_type = type(shape_model)

            if model_type is PptxPictureBoxModel:
                self.add_picture(slide, shape_model)

            elif model_type is PptxAutoShapeBoxModel:
                self.add_autoshape(slide, shape_model)

            elif model_type is PptxTextBoxModel:
                self.add_textbox(slide, shape_model)

            elif model_type is PptxConnectorModel:
                self.add_connector(slide, shape_model)

    def add_connector(self, slide: Slide, connector_model: PptxConnectorModel):
        if connector_model.thickness == 0:
            return
        connector_shape = slide.shapes.add_connector(
            connector_model.type, *connector_model.position.to_pt_xyxy()
        )
        connector_shape.line.width = Pt(connector_model.thickness)
        connector_shape.line.color.rgb = RGBColor.from_string(connector_model.color)
        self.set_fill_opacity(connector_shape, connector_model.opacity)

    def add_picture(self, slide: Slide, picture_model: PptxPictureBoxModel):
        image_path = picture_model.picture.path
        if (
            picture_model.clip
            or picture_model.border_radius
            or picture_model.invert
            or picture_model.opacity
            or picture_model.object_fit
            or picture_model.shape
        ):
            try:
                image = Image.open(image_path)
            except:
                print(f"Could not open image: {image_path}")
                return

            image = image.convert("RGBA")
            # ? Applying border radius twice to support both clip and object fit
            if picture_model.border_radius:
                image = round_image_corners(image, picture_model.border_radius)
            if picture_model.object_fit:
                image = fit_image(
                    image,
                    picture_model.position.width,
                    picture_model.position.height,
                    picture_model.object_fit,
                )
            elif picture_model.clip:
                image = clip_image(
                    image,
                    picture_model.position.width,
                    picture_model.position.height,
                )
            if picture_model.border_radius:
                image = round_image_corners(image, picture_model.border_radius)
            if picture_model.shape == PptxBoxShapeEnum.CIRCLE:
                image = create_circle_image(image)
            if picture_model.invert:
                image = invert_image(image)
            if picture_model.opacity:
                image = set_image_opacity(image, picture_model.opacity)
            image_path = os.path.join(self._temp_dir, f"{uuid.uuid4()}.png")
            image.save(image_path)

        margined_position = self.get_margined_position(
            picture_model.position, picture_model.margin
        )

        slide.shapes.add_picture(image_path, *margined_position.to_pt_list())

    def add_autoshape(self, slide: Slide, autoshape_box_model: PptxAutoShapeBoxModel):
        position = autoshape_box_model.position
        if autoshape_box_model.margin:
            position = self.get_margined_position(position, autoshape_box_model.margin)

        autoshape = slide.shapes.add_shape(
            autoshape_box_model.type, *position.to_pt_list()
        )

        textbox = autoshape.text_frame
        textbox.word_wrap = autoshape_box_model.text_wrap

        self.apply_fill_to_shape(autoshape, autoshape_box_model.fill)
        self.apply_margin_to_text_box(textbox, autoshape_box_model.margin)
        self.apply_stroke_to_shape(autoshape, autoshape_box_model.stroke)
        self.apply_shadow_to_shape(autoshape, autoshape_box_model.shadow)
        self.apply_border_radius_to_shape(autoshape, autoshape_box_model.border_radius)

        if autoshape_box_model.paragraphs:
            self.add_paragraphs(textbox, autoshape_box_model.paragraphs)

    def add_textbox(self, slide: Slide, textbox_model: PptxTextBoxModel):
        position = textbox_model.position
        textbox_shape = slide.shapes.add_textbox(*position.to_pt_list())
        textbox_shape.width += Pt(2)

        textbox = textbox_shape.text_frame
        textbox.word_wrap = textbox_model.text_wrap

        self.apply_fill_to_shape(textbox_shape, textbox_model.fill)
        self.apply_margin_to_text_box(textbox, textbox_model.margin)
        self.add_paragraphs(textbox, textbox_model.paragraphs)

    def add_paragraphs(
        self, textbox: TextFrame, paragraph_models: List[PptxParagraphModel]
    ):
        for index, paragraph_model in enumerate(paragraph_models):
            paragraph = textbox.add_paragraph() if index > 0 else textbox.paragraphs[0]
            self.populate_paragraph(paragraph, paragraph_model)

    def populate_paragraph(
        self, paragraph: _Paragraph, paragraph_model: PptxParagraphModel
    ):
        if paragraph_model.spacing:
            self.apply_spacing_to_paragraph(paragraph, paragraph_model.spacing)

        if paragraph_model.line_height:
            paragraph.line_spacing = paragraph_model.line_height

        if paragraph_model.alignment:
            paragraph.alignment = paragraph_model.alignment

        if paragraph_model.font:
            self.apply_font_to_paragraph(paragraph, paragraph_model.font)

        text_runs = []
        if paragraph_model.text:
            text_runs = self.parse_html_text_to_text_runs(
                paragraph_model.font, paragraph_model.text
            )
        elif paragraph_model.text_runs:
            text_runs = paragraph_model.text_runs

        for text_run_model in text_runs:
            text_run = paragraph.add_run()
            self.populate_text_run(text_run, text_run_model)

    def parse_html_text_to_text_runs(self, font: Optional[PptxFontModel], text: str):
        return parse_inline_html_to_runs(text, font)

    def populate_text_run(self, text_run: _Run, text_run_model: PptxTextRunModel):
        text_run.text = text_run_model.text
        if text_run_model.font:
            self.apply_font(text_run.font, text_run_model.font)

    def apply_border_radius_to_shape(self, shape: Shape, border_radius: Optional[int]):
        if not border_radius:
            return
        try:
            normalized_border_radius = Pt(border_radius) / min(
                shape.width, shape.height
            )
            shape.adjustments[0] = normalized_border_radius
        except:
            print("Could not apply border radius.")

    def apply_fill_to_shape(self, shape: Shape, fill: Optional[PptxFillModel] = None):
        if not fill:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(fill.color)
            self.set_fill_opacity(shape.fill, fill.opacity)

    def apply_stroke_to_shape(
        self, shape: Shape, stroke: Optional[PptxStrokeModel] = None
    ):
        if not stroke or stroke.thickness == 0:
            shape.line.fill.background()
        else:
            shape.line.fill.solid()
            shape.line.fill.fore_color.rgb = RGBColor.from_string(stroke.color)
            shape.line.width = Pt(stroke.thickness)
            self.set_fill_opacity(shape.line.fill, stroke.opacity)

    def apply_shadow_to_shape(
        self, shape: Shape, shadow: Optional[PptxShadowModel] = None
    ):

        # Access the XML for the shape
        sp_element = shape._element
        sp_pr = sp_element.xpath("p:spPr")[0]  # Shape properties XML element

        nsmap = sp_pr.nsmap

        # # Remove existing shadow effects if present
        effect_list = sp_pr.find("a:effectLst", namespaces=nsmap)
        if effect_list:
            old_outer_shadow = effect_list.find("a:outerShdw")
            if old_outer_shadow:
                effect_list.remove(
                    old_outer_shadow, namespaces=nsmap
                )  # Remove the old shadow
            old_inner_shadow = effect_list.find("a:innerShdw")
            if old_inner_shadow:
                effect_list.remove(
                    old_inner_shadow, namespaces=nsmap
                )  # Remove the old shadow
            old_prst_shadow = effect_list.find("a:prstShdw")
            if old_prst_shadow:
                effect_list.remove(
                    old_prst_shadow, namespaces=nsmap
                )  # Remove the old shadow

        if not effect_list:
            effect_list = etree.SubElement(
                sp_pr, f"{{{nsmap['a']}}}effectLst", nsmap=nsmap
            )

        if shadow is None:
            # Apply shadow with zero values when shadow is None
            outer_shadow = etree.SubElement(
                effect_list,
                f"{{{nsmap['a']}}}outerShdw",
                {
                    "blurRad": "0",
                    "dist": "0",
                    "dir": "0",
                },
                nsmap=nsmap,
            )
            color_element = etree.SubElement(
                outer_shadow,
                f"{{{nsmap['a']}}}srgbClr",
                {"val": "000000"},
                nsmap=nsmap,
            )
            etree.SubElement(
                color_element,
                f"{{{nsmap['a']}}}alpha",
                {"val": "0"},
                nsmap=nsmap,
            )
        else:
            # Apply the provided shadow
            outer_shadow = etree.SubElement(
                effect_list,
                f"{{{nsmap['a']}}}outerShdw",
                {
                    "blurRad": f"{Pt(shadow.radius)}",
                    "dir": f"{shadow.angle * 1000}",
                    "dist": f"{Pt(shadow.offset)}",
                    "rotWithShape": "0",
                },
                nsmap=nsmap,
            )
            color_element = etree.SubElement(
                outer_shadow,
                f"{{{nsmap['a']}}}srgbClr",
                {"val": f"{shadow.color}"},
                nsmap=nsmap,
            )
            etree.SubElement(
                color_element,
                f"{{{nsmap['a']}}}alpha",
                {"val": f"{int(shadow.opacity * 100000)}"},
                nsmap=nsmap,
            )

    def set_fill_opacity(self, fill, opacity):
        if opacity is None or opacity >= 1.0:
            return

        alpha = int((opacity) * 100000)

        try:
            ts = fill._xPr.solidFill
            sF = ts.get_or_change_to_srgbClr()
            self.get_sub_element(sF, "a:alpha", val=str(alpha))
        except Exception as e:
            print(f"Could not set fill opacity: {e}")

    def get_margined_position(
        self, position: PptxPositionModel, margin: Optional[PptxSpacingModel]
    ) -> PptxPositionModel:
        if not margin:
            return position

        left = position.left + margin.left
        top = position.top + margin.top
        width = max(position.width - margin.left - margin.right, 0)
        height = max(position.height - margin.top - margin.bottom, 0)

        return PptxPositionModel(left=left, top=top, width=width, height=height)

    def apply_margin_to_text_box(
        self, text_frame: TextFrame, margin: Optional[PptxSpacingModel]
    ) -> PptxPositionModel:
        text_frame.margin_left = Pt(margin.left if margin else 0)
        text_frame.margin_right = Pt(margin.right if margin else 0)
        text_frame.margin_top = Pt(margin.top if margin else 0)
        text_frame.margin_bottom = Pt(margin.bottom if margin else 0)

    def apply_spacing_to_paragraph(
        self, paragraph: _Paragraph, spacing: PptxSpacingModel
    ):
        paragraph.space_before = Pt(spacing.top)
        paragraph.space_after = Pt(spacing.bottom)

    def apply_font_to_paragraph(self, paragraph: _Paragraph, font: PptxFontModel):
        self.apply_font(paragraph.font, font)

    def apply_font(self, font: Font, font_model: PptxFontModel):
        font.name = font_model.name
        font.color.rgb = RGBColor.from_string(font_model.color)
        font.italic = font_model.italic
        font.size = Pt(font_model.size)
        font.bold = font_model.font_weight >= 600
        if font_model.underline is not None:
            font.underline = bool(font_model.underline)
        if font_model.strike is not None:
            self.apply_strike_to_font(font, font_model.strike)

    def apply_strike_to_font(self, font: Font, strike: Optional[bool]):
        try:
            rPr = font._element
            if strike is True:
                rPr.set("strike", "sngStrike")
            elif strike is False:
                rPr.set("strike", "noStrike")
        except Exception as e:
            print(f"Could not apply strikethrough: {e}")

    def save(self, path: str):
        self._ppt.save(path)
